"""
GeoCharge Montréal — Application Python Shiny
GMQ580 Géomatique Informatique 2 — Été 2026 — Université de Sherbrooke

Remplace Django+Docker+PostGIS par une app autonome :
  - geopandas : requêtes spatiales (sjoin, buffer, within)
  - folium    : carte Leaflet interactive
  - shiny     : interface réactive

Lancer : shiny run app.py --reload
"""

import sys, os, math, warnings, unicodedata, json, csv, io
from datetime import datetime

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

import pandas as pd
import geopandas as gpd
import docx
from docx.shared import Pt, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptxIn, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import folium
from folium import plugins
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from shiny import App, ui, render, reactive

warnings.filterwarnings("ignore", category=FutureWarning)

# ═══════════════════════════════════════════════════════════════════════════
# CHEMINS DES DONNÉES
# ═══════════════════════════════════════════════════════════════════════════
_BASE = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
VEC  = os.path.join(_BASE, "data", "vectors")
DATA = os.path.join(_BASE, "data")

# ── Journal d'activité (persistant, sur disque) ────────────────────────────
# Chaque recherche sur la carte et chaque analyse gestionnaire (G1-G5) est
# ajoutée à ce fichier CSV, qui survit à la fermeture de l'app — c'est la base
# des 3 rapports offerts dans l'onglet « Rapports » (historique de recherche,
# résumé des analyses, journal quotidien complet).
JOURNAL_PATH = os.path.join(os.path.dirname(__file__), "journal_activite.csv")
JOURNAL_COLONNES = ["date", "heure", "categorie", "detail"]

LOGO_PATH = os.path.join(_BASE, "UdeS.jpg")
LOGO_DATA_URI = None
if os.path.exists(LOGO_PATH):
    import base64 as _b64
    with open(LOGO_PATH, "rb") as _f:
        LOGO_DATA_URI = "data:image/jpeg;base64," + _b64.b64encode(_f.read()).decode("ascii")


def _init_journal():
    if not os.path.exists(JOURNAL_PATH):
        pd.DataFrame(columns=JOURNAL_COLONNES).to_csv(JOURNAL_PATH, index=False)


def _log_event(categorie, detail):
    now = datetime.now()
    row = {"date": now.strftime("%Y-%m-%d"), "heure": now.strftime("%H:%M:%S"),
           "categorie": categorie, "detail": detail}
    file_existe = os.path.exists(JOURNAL_PATH)
    with open(JOURNAL_PATH, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=JOURNAL_COLONNES)
        if not file_existe:
            writer.writeheader()
        writer.writerow(row)


def _lire_journal():
    _init_journal()
    try:
        df = pd.read_csv(JOURNAL_PATH, dtype=str)
        for col in JOURNAL_COLONNES:
            if col not in df.columns:
                df[col] = ""
        return df[JOURNAL_COLONNES]
    except pd.errors.EmptyDataError:
        return pd.DataFrame(columns=JOURNAL_COLONNES)


_init_journal()


# ═══════════════════════════════════════════════════════════════════════════
# CHARGEMENT ET PRÉ-CALCUL (exécuté une fois au démarrage)
# ═══════════════════════════════════════════════════════════════════════════
print("⏳ Chargement des données GeoJSON…", flush=True)

bornes_gdf    = gpd.read_file(os.path.join(VEC, "bornes_recharge_montreal.geojson")).to_crs(4326)
arrond_gdf    = gpd.read_file(os.path.join(VEC, "arrondissements_montreal.geojson")).to_crs(4326)
parcs_gdf     = gpd.read_file(os.path.join(VEC, "parcs_montreal.geojson")).to_crs(4326)
epiceries_gdf = gpd.read_file(os.path.join(VEC, "epiceries_montreal.geojson")).to_crs(4326)
stations_gdf  = gpd.read_file(os.path.join(VEC, "stations_metro_stm.geojson")).to_crs(4326)
demo_df       = pd.read_csv(os.path.join(DATA, "demo_arrondissements.csv"))

# Normaliser la colonne nom dans arrond_gdf
for col in arrond_gdf.columns:
    if col.upper() == "NOM" and col != "nom":
        arrond_gdf = arrond_gdf.rename(columns={col: "nom"})
        break

# demo_arrondissements.csv utilise le tiret cadratin (–) dans les noms composés
# (ex. "Côte-des-Neiges–Notre-Dame-de-Grâce") alors que le GeoJSON des
# arrondissements utilise le trait d'union simple (-) partout : sans cette
# normalisation, la jointure échoue silencieusement sur 6 des 19 arrondissements.
demo_df["arrondissement"] = demo_df["arrondissement"].str.replace("–", "-", regex=False)

print(f"  → {len(bornes_gdf)} bornes · {len(arrond_gdf)} arrondissements · "
      f"{len(parcs_gdf)} parcs · {len(epiceries_gdf)} épiceries · "
      f"{len(stations_gdf)} stations métro", flush=True)

# ── Projeter en MTM-8 (EPSG:32188) pour les calculs de distance ────────────
bornes_p    = bornes_gdf.to_crs(32188)
arrond_p    = arrond_gdf.to_crs(32188)
parcs_p     = parcs_gdf.to_crs(32188)
epiceries_p = epiceries_gdf.to_crs(32188)
stations_p  = stations_gdf.to_crs(32188)

# ── Couverture 500 m par arrondissement ────────────────────────────────────
print("⏳ Calcul couverture 500 m…", flush=True)
buf_union = bornes_p.geometry.buffer(500).unary_union
arrond_p["pct_couverture"] = arrond_p.geometry.apply(
    lambda g: round(min(100.0, 100.0 * buf_union.intersection(g).area / g.area), 1)
)
arrond_gdf["pct_couverture"] = arrond_p["pct_couverture"].values

# ── Nombre de bornes par arrondissement ────────────────────────────────────
_b_in_a = gpd.sjoin(bornes_p[["geometry"]], arrond_p[["nom", "geometry"]],
                    how="left", predicate="within")
_nb = _b_in_a.groupby("nom").size().reset_index(name="nb_bornes")
arrond_gdf = arrond_gdf.merge(_nb, on="nom", how="left")
arrond_gdf["nb_bornes"] = arrond_gdf["nb_bornes"].fillna(0).astype(int)

# ── Fusion données socio-démographiques ────────────────────────────────────
demo_df = demo_df.rename(columns={"arrondissement": "nom"})
arrond_gdf = arrond_gdf.merge(
    demo_df[["nom", "pop_2021", "densite_pop_km2", "revenu_median_menage",
             "tx_voiture_pct", "tx_faible_revenu_pct"]],
    on="nom", how="left"
)

# ── Score de priorité composite (5 critères pondérés, cf. RAPPORT_FINAL §3.7) ─
# Restreint aux 19 arrondissements avec données StatCan (les 15 villes liées
# n'ont pas de profil démographique officiel par arrondissement).
_score_df = arrond_gdf.dropna(subset=["densite_pop_km2", "revenu_median_menage",
                                       "tx_voiture_pct", "tx_faible_revenu_pct"]).copy()
_gap    = 1 - _score_df["pct_couverture"] / 100
_dens   = _score_df["densite_pop_km2"] / _score_df["densite_pop_km2"].max()
_moto   = _score_df["tx_voiture_pct"] / _score_df["tx_voiture_pct"].max()
_equite = 1 - _score_df["revenu_median_menage"] / _score_df["revenu_median_menage"].max()
_fr     = _score_df["tx_faible_revenu_pct"] / _score_df["tx_faible_revenu_pct"].max()

_score_df["score_priorite"] = (100 * (
    _gap * 0.35 + _dens * 0.25 + _moto * 0.15 + _equite * 0.15 + _fr * 0.10
)).round(1)

priorite_df = _score_df[["nom", "pct_couverture", "nb_bornes", "score_priorite"]] \
    .sort_values("score_priorite", ascending=False).reset_index(drop=True)

# ── Bornes à 500 m pour chaque parc ────────────────────────────────────────
print("⏳ Requête spatiale parcs ↔ bornes (500 m)…", flush=True)
_buf_bornes_500 = bornes_p.copy()
_buf_bornes_500["geometry"] = _buf_bornes_500.geometry.buffer(500)
_buf_bornes_500["_bid"] = range(len(_buf_bornes_500))
_pj = gpd.sjoin(parcs_p[["geometry"]], _buf_bornes_500[["geometry", "_bid"]],
                how="left", predicate="intersects")
parcs_gdf["nb_bornes_500m"] = (
    _pj.groupby(level=0)["_bid"].count().reindex(parcs_p.index, fill_value=0).values
)

# ── Épiceries sans borne à 300 m ───────────────────────────────────────────
print("⏳ Requête spatiale épiceries ↔ bornes (300 m)…", flush=True)
_buf_bornes_300_union = bornes_p.geometry.buffer(300).unary_union
epiceries_gdf["has_borne_300m"] = epiceries_p.geometry.within(_buf_bornes_300_union)

# ── Parcs et épiceries par arrondissement (pour G6) ───────────────────────
_pa = gpd.sjoin(parcs_p[["geometry"]], arrond_p[["nom", "geometry"]],
                how="left", predicate="within")
_nb_parcs = _pa.groupby("nom").size().reset_index(name="nb_parcs")

_ea = gpd.sjoin(epiceries_p[["geometry"]], arrond_p[["nom", "geometry"]],
                how="left", predicate="within")
_nb_epic = _ea.groupby("nom").size().reset_index(name="nb_epiceries")

g6_df = arrond_gdf[["nom", "nb_bornes", "pct_couverture",
                     "densite_pop_km2", "revenu_median_menage",
                     "tx_voiture_pct", "tx_faible_revenu_pct"]].copy()
g6_df = g6_df.merge(_nb_parcs, on="nom", how="left")
g6_df = g6_df.merge(_nb_epic,  on="nom", how="left")
g6_df["nb_parcs"]     = g6_df["nb_parcs"].fillna(0).astype(int)
g6_df["nb_epiceries"] = g6_df["nb_epiceries"].fillna(0).astype(int)

print("✅ Données prêtes.", flush=True)


# ═══════════════════════════════════════════════════════════════════════════
# UTILITAIRES
# ═══════════════════════════════════════════════════════════════════════════

def pearson(xs, ys):
    n = len(xs)
    if n < 3:
        return 0.0
    mx, my = sum(xs) / n, sum(ys) / n
    num = sum((x - mx) * (y - my) for x, y in zip(xs, ys))
    dx  = sum((x - mx) ** 2 for x in xs) ** 0.5
    dy  = sum((y - my) ** 2 for y in ys) ** 0.5
    return round(num / (dx * dy + 1e-9), 3)


def coul_couverture(pct):
    if pct >= 60: return "#00c96e"
    if pct >= 30: return "#ffd700"
    if pct >= 15: return "#ff8c00"
    return "#e8273e"


DEMO_LABELS = {
    "revenu_median_menage": "Revenu médian des ménages ($)",
    "tx_voiture_pct":       "Taux de motorisation (%)",
    "densite_pop_km2":      "Densité de population (hab/km²)",
    "tx_faible_revenu_pct": "Faible revenu (% ménages < 40 k$)",
}

DEMO_PALETTES = {
    "revenu_median_menage": "RdYlGn",
    "tx_voiture_pct":       "YlGnBu",
    "densite_pop_km2":      "Blues",
    "tx_faible_revenu_pct": "YlOrRd",
}

# ── Pré-calculs pour la carte ───────────────────────────────────────────────
# build_map() est rappelée à chaque interaction (calque, fond de carte, variable
# démographique) : on précalcule ici tout ce qui ne dépend pas de ces choix, pour
# ne pas refaire une boucle Python sur 2 412 bornes ni resérialiser le GeoJSON des
# arrondissements à chaque clic.
BORNES_LATLON = list(zip(bornes_gdf.geometry.y, bornes_gdf.geometry.x))

_ARROND_COUVERTURE_JSON = arrond_gdf[["nom", "nb_bornes", "pct_couverture", "geometry"]].to_json()

_ARROND_DEMO_VALID = {
    var: arrond_gdf[["nom", var, "geometry"]].dropna(subset=[var])
    for var in DEMO_LABELS
}
_ARROND_DEMO_JSON = {var: gdf.to_json() for var, gdf in _ARROND_DEMO_VALID.items()}

# ── Barre de recherche personnalisée (tolérante aux accents, tirets, mots
# manquants et abréviations "St-"/"Ste-") ───────────────────────────────────
# Le plugin Leaflet Search standard ne fait qu'une comparaison de sous-chaîne
# stricte : "cote des neiges", "ile-bizard" ou "St-Laurent" ne trouvaient rien.
# On construit ici, pour chaque arrondissement, une clé normalisée + son
# emprise (bounds), utilisées par un petit script JS qui compare mot par mot.
def _clef_recherche(texte):
    sans_accents = "".join(
        c for c in unicodedata.normalize("NFD", texte) if unicodedata.category(c) != "Mn"
    )
    sans_ponctuation = sans_accents.replace("-", " ").replace("'", " ").replace("’", " ")
    return " ".join(sans_ponctuation.lower().split())

_arrond_bounds = arrond_gdf.geometry.bounds  # minx, miny, maxx, maxy
_ARROND_SEARCH_DATA = [
    {
        "nom": nom,
        "clef": _clef_recherche(nom),
        "bounds": [[b.miny, b.minx], [b.maxy, b.maxx]],
    }
    for nom, b in zip(arrond_gdf["nom"], _arrond_bounds.itertuples())
]

VAR_LABELS_G6 = {
    "nb_parcs":             "Zones de loisirs (parcs)",
    "nb_epiceries":         "Commerces alimentaires (épiceries)",
    "revenu_median_menage": "Revenu médian des ménages",
    "densite_pop_km2":      "Densité de population",
    "tx_voiture_pct":       "Taux de motorisation (proxy demande EV)",
    "tx_faible_revenu_pct": "Taux de faible revenu",
}

QUESTIONS_G6 = {
    "nb_parcs":             "est-ce lié à la proximité de zones de loisirs ?",
    "nb_epiceries":         "est-ce la présence de commerces ?",
    "revenu_median_menage": "est-ce lié au profil économique ?",
    "densite_pop_km2":      "est-ce lié à la densité de population ?",
    "tx_voiture_pct":       "est-ce lié à l'accès en voiture ?",
    "tx_faible_revenu_pct": "est-ce lié à la vulnérabilité socio-éco ?",
}


# ═══════════════════════════════════════════════════════════════════════════
# CONSTRUCTEUR DE CARTE FOLIUM
# ═══════════════════════════════════════════════════════════════════════════

def build_map(show_bornes=True, show_arrond=True, demo_var="couverture"):
    m = folium.Map(location=[45.53, -73.62], zoom_start=11, tiles=None,
                   prefer_canvas=True)

    # ── Fonds de carte (un seul actif à la fois, changeable dans le contrôle) ──
    folium.TileLayer(
        "https://{s}.basemaps.cartocdn.com/dark_all/{z}/{x}/{y}{r}.png",
        attr="© OpenStreetMap, © CARTO", name="Plan sombre",
        max_zoom=19, show=True,
    ).add_to(m)
    folium.TileLayer(
        "https://{s}.basemaps.cartocdn.com/light_all/{z}/{x}/{y}{r}.png",
        attr="© OpenStreetMap, © CARTO", name="Plan clair",
        max_zoom=19, show=False,
    ).add_to(m)
    folium.TileLayer(
        "https://server.arcgisonline.com/ArcGIS/rest/services/"
        "World_Imagery/MapServer/tile/{z}/{y}/{x}",
        attr="© Esri, Maxar, Earthstar Geographics", name="Satellite",
        max_zoom=19, show=False,
    ).add_to(m)
    folium.TileLayer(
        "OpenStreetMap", name="OpenStreetMap", max_zoom=19, show=False,
    ).add_to(m)

    # ── Arrondissements ────────────────────────────────────────────────────
    if show_arrond:
        if demo_var in DEMO_LABELS:
            # Choroplèthe démographique (données précalculées au démarrage)
            _col = demo_var
            _valid_json = _ARROND_DEMO_JSON[_col]
            folium.Choropleth(
                geo_data=_valid_json,
                data=_ARROND_DEMO_VALID[_col],
                columns=["nom", _col],
                key_on="feature.properties.nom",
                fill_color=DEMO_PALETTES.get(_col, "YlOrRd"),
                fill_opacity=0.65,
                line_opacity=0.8,
                line_color="#1e3a5f",
                nan_fill_color="#222",
                legend_name=DEMO_LABELS[_col],
                name="Arrondissements (démographie)",
                control=False,  # déjà contrôlé par la case à cocher du panneau de gauche
            ).add_to(m)
            # Tooltip overlay
            folium.GeoJson(
                _valid_json,
                style_function=lambda x: {"fillOpacity": 0, "weight": 0},
                tooltip=folium.GeoJsonTooltip(
                    fields=["nom", _col],
                    aliases=["Arrondissement", DEMO_LABELS[_col]],
                    localize=True,
                ),
                name="_tooltip_demo",
                control=False,  # couche technique (info-bulles) — ne doit pas apparaître dans le contrôle des calques
            ).add_to(m)
        else:
            # Couverture en bornes
            def _style(feat):
                pct = feat["properties"].get("pct_couverture") or 0
                return {
                    "fillColor": coul_couverture(pct),
                    "color": "#1e3a5f",
                    "weight": 1.2,
                    "fillOpacity": 0.45,
                }
            folium.GeoJson(
                _ARROND_COUVERTURE_JSON,
                name="Arrondissements (couverture)",
                style_function=_style,
                tooltip=folium.GeoJsonTooltip(
                    fields=["nom", "nb_bornes", "pct_couverture"],
                    aliases=["Arrondissement", "Bornes", "Couverture (%)"],
                    localize=True,
                ),
                control=False,  # déjà contrôlé par la case à cocher du panneau de gauche
            ).add_to(m)

    # ── Bornes de recharge ─────────────────────────────────────────────────
    if show_bornes:
        plugins.FastMarkerCluster(
            BORNES_LATLON,
            name="Bornes de recharge",
            options={"maxClusterRadius": 40, "disableClusteringAtZoom": 15},
            control=False,  # déjà contrôlé par la case à cocher du panneau de gauche
        ).add_to(m)

    folium.LayerControl(collapsed=False).add_to(m)

    # ── Point d'entrée pour la barre de recherche du panneau latéral (voir
    # SIDEBAR_SEARCH_JS) : la recherche elle-même vit hors de cette iframe, dans la page
    # Shiny ; on expose ici juste une fonction que ce code externe peut appeler pour
    # centrer/zoomer la carte et marquer l'arrondissement trouvé.
    _map_var = m.get_name()
    _zoom_js = f"""
    (function() {{
      var marker = null;
      window.geochargeZoomTo = function(bounds) {{
        {_map_var}.fitBounds(bounds);
        if (marker) {_map_var}.removeLayer(marker);
        var c = [(bounds[0][0]+bounds[1][0])/2, (bounds[0][1]+bounds[1][1])/2];
        marker = L.circleMarker(c, {{radius:14, color:'#00e5ff', weight:3, fillOpacity:0}}).addTo({_map_var});
      }};
    }})();
    """
    m.get_root().script.add_child(folium.Element(_zoom_js))

    # ── Habillage néon des contrôles Leaflet (dans le document de l'iframe) ─
    m.get_root().header.add_child(folium.Element("""
    <style>
      .leaflet-control-layers, .leaflet-bar {
        background: #0d1b2e !important;
        border: 1px solid rgba(0,229,255,0.45) !important;
        box-shadow: 0 0 14px -2px rgba(0,229,255,0.5) !important;
        color: #d9eef8 !important;
      }
      .leaflet-control-layers label, .leaflet-control-layers span { color: #d9eef8 !important; }
      .leaflet-bar a { background: #0d1b2e !important; color: #00e5ff !important; border-color: rgba(0,229,255,0.3) !important; }
      .leaflet-bar a:hover { background: #102540 !important; box-shadow: 0 0 8px rgba(0,229,255,0.6) inset; }
      .leaflet-control-layers input[type="radio"] { accent-color: #00e5ff; }
    </style>
    """))
    return m


# ═══════════════════════════════════════════════════════════════════════════
# INTERFACE UTILISATEUR SHINY
# ═══════════════════════════════════════════════════════════════════════════

NEON_CSS = """
:root {
  --bg-deep: #060b14;
  --bg-panel: #0d1b2e;
  --bg-panel-2: #102540;
  --neon: #00e5ff;
  --neon-soft: rgba(0,229,255,0.35);
  --neon-glow: 0 0 6px rgba(0,229,255,0.85), 0 0 18px rgba(0,229,255,0.45), 0 0 34px rgba(0,229,255,0.18);
  --ink: #d9eef8;
  --ink-dim: #82a8c0;
}
html, body { background: var(--bg-deep) !important; color: var(--ink) !important; }
.navbar { background: var(--bg-panel) !important; border-bottom: 1px solid var(--neon-soft) !important;
  box-shadow: 0 0 22px rgba(0,229,255,0.28); }
.navbar-brand span, .navbar-brand { text-shadow: 0 0 8px rgba(0,229,255,0.75), 0 0 20px rgba(0,229,255,0.35) !important; }
/* Petits cadrages : le titre et chaque onglet ressortent comme des blocs distincts plutôt
   que du texte flottant dans la barre. */
.navbar-brand {
  border: 1px solid rgba(0,229,255,0.3); border-radius: 6px;
  padding: 4px 12px !important; margin-right: 10px !important;
  background: rgba(0,229,255,0.05);
}
/* Étale les onglets sur toute la largeur restante de la barre plutôt que de les laisser
   groupés à gauche. */
.navbar .container-fluid { align-items: center; }
.navbar-collapse { flex-grow: 1; }
.navbar-nav { gap: 6px; flex-grow: 1; justify-content: space-evenly; }
.nav-item { flex: 1; display: flex; }
.nav-link {
  color: var(--ink-dim) !important;
  border: 1px solid rgba(0,229,255,0.18); border-radius: 6px;
  padding: 6px 14px !important; margin: 4px 0;
  flex: 1; text-align: center; justify-content: center;
}
.nav-link.active {
  color: var(--neon) !important; text-shadow: 0 0 8px rgba(0,229,255,0.6);
  border-color: rgba(0,229,255,0.55); background: rgba(0,229,255,0.08);
  box-shadow: 0 0 10px -3px rgba(0,229,255,0.6);
}
.nav-link:hover { border-color: rgba(0,229,255,0.4); }

.bslib-sidebar-layout > div.sidebar, .sidebar {
  background: var(--bg-panel) !important; border-right: 1px solid var(--neon-soft) !important;
  box-shadow: inset -10px 0 24px -20px var(--neon);
  color: var(--ink) !important;
}
.sidebar h6 { color: var(--neon) !important; text-shadow: 0 0 6px rgba(0,229,255,0.5); letter-spacing: 0.04em; }
.sidebar label, .sidebar p, .sidebar strong, .sidebar span, .sidebar div { color: var(--ink) !important; }
.sidebar p[style*="666"], .sidebar p[style*="999"] { color: var(--ink-dim) !important; }
/* La marge par défaut des <hr> (Bootstrap) s'ajoutait au "gap" déjà fixe entre les contrôles
   du panneau, doublant l'espace visible autour de chaque séparateur. On resserre les deux. */
.sidebar-content { gap: 14px !important; }
.sidebar-content hr { margin: 0 !important; }
/* Accent néon vertical clignotant (comme un curseur) sur le bord gauche de la barre de
   recherche. Élément séparé plutôt qu'un ::before sur l'<input> (mal supporté sur les champs
   de formulaire selon les navigateurs). */
#sidebar-search-accent {
  position: absolute; left: -3px; top: 3px; bottom: 3px; width: 3px;
  background: var(--neon); border-radius: 2px;
  box-shadow: 0 0 8px 2px rgba(0,229,255,0.85);
  animation: neon-blink 1.1s steps(1, start) infinite;
}
@keyframes neon-blink {
  0%, 49%   { opacity: 1; }
  50%, 100% { opacity: 0.12; }
}

.card { background: var(--bg-panel) !important; border: 1px solid var(--neon-soft) !important;
  box-shadow: 0 0 16px -6px rgba(0,229,255,0.35); }
.card-header { background: var(--bg-panel-2) !important; color: var(--neon) !important;
  border-bottom: 1px solid var(--neon-soft) !important; text-shadow: 0 0 6px rgba(0,229,255,0.45); font-weight: 700; }
.card p, .card label { color: var(--ink-dim) !important; }

table { color: var(--ink) !important; }
.table > :not(caption) > * > * { background-color: transparent !important; color: var(--ink) !important;
  border-color: rgba(0,229,255,0.15) !important; }
thead { color: var(--neon) !important; }

select, input[type=number], input[type=text], .form-control, .form-select {
  background: var(--bg-deep) !important; color: var(--ink) !important;
  border: 1px solid var(--neon-soft) !important;
}
select:focus, input:focus { outline: none !important; box-shadow: var(--neon-glow) !important; border-color: var(--neon) !important; }

/* Shiny transforme les <select> en Selectize.js, qui tronque le texte affiché
   par défaut (ex: "Couverture en bornes" coupé en "Couverture en born") — on
   force l'affichage complet, quitte à ce que la ligne s'agrandisse un peu. */
.selectize-input, .selectize-input.items, .selectize-input .item {
  white-space: normal !important;
  overflow: visible !important;
  text-overflow: clip !important;
  background: var(--bg-deep) !important;
  color: var(--ink) !important;
  border-color: rgba(0,229,255,0.45) !important;
}
.selectize-dropdown, .selectize-dropdown .option {
  background: var(--bg-panel) !important;
  color: var(--ink) !important;
}
.selectize-dropdown .active { background: var(--neon) !important; color: #06222c !important; }
/* Selectize ajoute un champ texte caché pour filtrer les options ; son curseur clignotant
   reste visible même si l'on ne tape pas de recherche — on le masque. */
.selectize-input input { caret-color: transparent !important; }

.btn-primary { background: linear-gradient(180deg, #00c2e6, #0089a8) !important; border: none !important;
  box-shadow: 0 0 10px rgba(0,229,255,0.55); font-weight: 700; }
.btn-primary:hover { box-shadow: var(--neon-glow) !important; transform: translateY(-1px); }

.alert { border: 1px solid rgba(255,255,255,0.18) !important; box-shadow: 0 0 12px rgba(0,0,0,0.45); font-weight: 600; }

footer, .navbar + div { color: var(--ink-dim) !important; }

/* Cadre néon autour de la carte interactive */
.html-fill-item:has(iframe), iframe {
  border: 1.5px solid var(--neon) !important;
  border-radius: 6px;
  box-shadow: var(--neon-glow);
}

/* Resserre l'espacement par défaut de bslib autour de la ligne carte+panneau (padding du
   conteneur "main" + marge sous la ligne), qui laissait un vide voyant entre la carte et le
   pied de page. ADJUST_MAP_HEIGHT_JS mesure ces valeurs dynamiquement, donc la carte
   récupère automatiquement l'espace ainsi libéré. */
.bslib-sidebar-layout .main { padding-top: 6px !important; padding-bottom: 6px !important; }
.bslib-sidebar-layout { margin-bottom: 6px !important; }
"""

# Ajuste la hauteur de la carte à l'espace réellement disponible (fenêtre - navbar - pied de
# page), plutôt qu'un calc(100vh - Xpx) figé qui déborde ou laisse un vide selon la fenêtre.
ADJUST_MAP_HEIGHT_JS = """
function ajusterHauteurCarte() {
  var frame = document.getElementById('geocharge-map-frame');
  var row = document.querySelector('.bslib-sidebar-layout');
  if (!frame || !row) return;
  // Chrome (padding/marge/bordure) entre la carte et la ligne carte+panneau qui la contient
  // (ex. le padding du conteneur "main" de bslib) — valeurs CSS fixes, indépendantes du
  // contenu, donc mesurées une fois par simple somme plutôt que devinées.
  var chromeBelow = 0;
  for (var el = frame; el && el !== row; el = el.parentElement) {
    var cs = getComputedStyle(el);
    chromeBelow += parseFloat(cs.paddingBottom) || 0;
    chromeBelow += parseFloat(cs.marginBottom) || 0;
    chromeBelow += parseFloat(cs.borderBottomWidth) || 0;
  }
  var frameTop = frame.getBoundingClientRect().top;
  // Sonde : le <body> a un min-height:100% (calque bslib), donc scrollHeight reste bloqué à
  // innerHeight tant que le contenu réel est plus court que la fenêtre — une mesure prise dans
  // cet état est fausse et, réappliquée en boucle (setInterval), fait s'effondrer la carte au
  // minimum (300px). On agrandit temporairement la carte pour garantir que le contenu dépasse
  // la fenêtre, on mesure alors le chrome après la ligne (marge + pied de page, invariant),
  // puis on applique la vraie hauteur cible en une seule fois (aucun scintillement : les deux
  // écritures de style se font avant le prochain rendu du navigateur).
  frame.style.height = (window.innerHeight + 800) + 'px';
  var belowRow = document.documentElement.scrollHeight
    - (window.scrollY + row.getBoundingClientRect().bottom);
  var h = Math.max(window.innerHeight - frameTop - belowRow - chromeBelow - 10, 300);
  frame.style.height = h + 'px';
  // Le panneau de gauche (checkboxes/menus) a une hauteur de contenu fixe qui, sans plafond,
  // peut forcer toute la ligne (donc la page) à s'agrandir au-delà de la fenêtre sur les petits
  // écrans. On ne le plafonne QUE si son contenu naturel dépasse réellement l'espace disponible
  // (sidebar.scrollHeight reflète toujours la taille naturelle, même sous max-height) — sinon on
  // le laisse à sa hauteur naturelle : sinon la mise en page flex (gap fixe) étire l'espace
  // excédentaire entre les contrôles au lieu de les garder groupés en haut.
  var sidebar = document.querySelector('.bslib-sidebar-layout .sidebar-content');
  if (sidebar) {
    if (sidebar.scrollHeight > h) {
      sidebar.style.maxHeight = h + 'px';
      sidebar.style.overflowY = 'auto';
    } else {
      sidebar.style.maxHeight = '';
      sidebar.style.overflowY = '';
    }
  }
}
window.addEventListener('resize', ajusterHauteurCarte);
window.addEventListener('load', ajusterHauteurCarte);
// Ce script s'exécute dans <head>, avant que document.body existe : observer/setTimeout ne
// peuvent être posés qu'une fois le DOM prêt, sinon MutationObserver.observe(document.body)
// lève une exception qui interrompt tout le reste du script (setInterval inclus).
document.addEventListener('DOMContentLoaded', function () {
  var obs = new MutationObserver(ajusterHauteurCarte);
  obs.observe(document.body, {childList: true, subtree: true});
});
setInterval(ajusterHauteurCarte, 500);
"""

# Barre de recherche d'arrondissement du panneau latéral. Vit dans la page Shiny (pas dans
# l'iframe folium comme avant) : pour centrer/zoomer la carte, elle appelle la fonction
# window.geochargeZoomTo exposée par build_map() à l'intérieur de l'iframe.
SIDEBAR_SEARCH_JS = f"""
(function() {{
  var DATA = {json.dumps(_ARROND_SEARCH_DATA, ensure_ascii=False)};
  function normalise(s) {{
    return s.normalize('NFD').replace(/[\\u0300-\\u036f]/g, '')
            .replace(/['’-]/g, ' ').toLowerCase().replace(/\\s+/g, ' ').trim();
  }}
  function expand(tok) {{
    if (tok === 'st') return 'saint';
    if (tok === 'ste') return 'sainte';
    return tok;
  }}
  function ready() {{
    var input = document.getElementById('sidebar-arr-search-input');
    var results = document.getElementById('sidebar-arr-search-results');
    if (!input || !results) return;
    input.addEventListener('input', function() {{
      var q = normalise(input.value);
      results.innerHTML = '';
      if (!q) return;
      var tokens = q.split(' ').filter(Boolean).map(expand);
      var matches = DATA.filter(function(d) {{
        var words = d.clef.split(' ');
        return tokens.every(function(t) {{
          return words.some(function(w) {{ return w.indexOf(t) === 0; }});
        }});
      }}).slice(0, 8);
      matches.forEach(function(d) {{
        var item = document.createElement('div');
        item.textContent = d.nom;
        item.style.cssText = 'padding:6px 10px; font-size:12.5px; cursor:pointer; ' +
          'background:#0d1b2e; color:#d9eef8; border:1px solid rgba(0,229,255,0.3); ' +
          'border-top:none;';
        item.onmouseover = function() {{ item.style.background = '#102540'; }};
        item.onmouseout  = function() {{ item.style.background = '#0d1b2e'; }};
        item.onclick = function() {{
          var frame = document.querySelector('#geocharge-map-frame iframe');
          try {{
            if (frame && frame.contentWindow && frame.contentWindow.geochargeZoomTo) {{
              frame.contentWindow.geochargeZoomTo(d.bounds);
            }}
          }} catch (e) {{}}
          results.innerHTML = '';
          input.value = d.nom;
          if (window.Shiny) {{
            Shiny.setInputValue('arr_search_click', {{nom: d.nom, ts: Date.now()}}, {{priority: 'event'}});
          }}
        }};
        results.appendChild(item);
      }});
    }});
  }}
  document.addEventListener('DOMContentLoaded', ready);
}})();
"""

# ui.update_date(value=None) ne vide pas réellement le composant bootstrap-datepicker sous-
# jacent (il ignore une valeur nulle) — on appelle donc directement sa méthode clearDates(),
# ce qui déclenche l'événement que la liaison Shiny du champ de date écoute pour renvoyer
# None au serveur.
RAPPORT_DATE_RESET_JS = """
document.addEventListener('click', function(e) {
  var link = e.target.closest && e.target.closest('#rapport-date-reset-link');
  if (!link) return;
  e.preventDefault();
  var input = document.querySelector('#rapport_date input');
  if (input && window.jQuery) {
    jQuery(input).bsDatepicker('clearDates');
  }
});
"""

# ── Bande de cases KPI sous la barre de navigation (une case par statistique, avec fond
# teinté propre à chacune) — visible sur tous les onglets. Mise en rangée séparée : les
# insérer dans la ligne de la navbar elle-même la faisait passer sur deux lignes une fois
# les onglets + 4 statistiques trop larges pour la fenêtre. ────────────────────────────
_KPI_HAUT = [
    (f"{len(bornes_gdf):,}".replace(",", " "), "Bornes de recharge"),
    (str(len(arrond_gdf)), "Arrondissements"),
    (f"{arrond_gdf['pct_couverture'].mean():.0f}%", "Couverture moyenne"),
    (str(len(stations_gdf)), "Stations de métro"),
]
# Grille compacte 2x2 (plutôt qu'une rangée de 4), pour tenir dans la largeur du panneau
# latéral — placée juste au-dessus de la barre de recherche.
KPI_BAR = ui.div(
    *[
        ui.div(
            ui.div(valeur, style="font-size:1.05rem; font-weight:700; color:var(--neon);"),
            ui.div(label, style="font-size:0.68rem; color:var(--ink-dim);"),
            style="text-align:center; padding:6px 4px; background:rgba(0,229,255,0.07); "
                  "border:1px solid rgba(0,229,255,0.25); border-radius:6px;",
        )
        for valeur, label in _KPI_HAUT
    ],
    style="display:grid; grid-template-columns:1fr 1fr; gap:8px; margin-bottom:14px;",
)

# Logo de l'Université de Sherbrooke dans le coin haut-droit de la barre de navigation
# (comme le logo institutionnel des tableaux de bord de référence), en plus du footer.
LOGO_NAVBAR_JS = f"""
document.addEventListener('DOMContentLoaded', function() {{
  var containerFluid = document.querySelector('.navbar .container-fluid');
  if (!containerFluid) return;
  var badge = document.createElement('span');
  badge.style.cssText = 'margin-left:auto; display:inline-flex; align-items:center; '
    + 'background:#ffffff; border-radius:4px; padding:3px 8px;';
  var img = document.createElement('img');
  img.src = {json.dumps(LOGO_DATA_URI or "")};
  img.style.cssText = 'height:24px; display:block;';
  badge.appendChild(img);
  containerFluid.appendChild(badge);
}});
"""

app_ui = ui.page_navbar(
    ui.head_content(
        ui.tags.style(NEON_CSS),
        ui.tags.script(ADJUST_MAP_HEIGHT_JS),
        ui.tags.script(RAPPORT_DATE_RESET_JS),
        ui.tags.script(SIDEBAR_SEARCH_JS),
    ),

    # ── TAB 1 : CARTE ──────────────────────────────────────────────────────
    ui.nav_panel(
        "🗺️ Carte interactive",
        ui.layout_sidebar(
            ui.sidebar(
                KPI_BAR,
                ui.div(
                    ui.div(id="sidebar-search-accent"),
                    ui.tags.input(
                        id="sidebar-arr-search-input", type="text", autocomplete="off",
                        placeholder="Rechercher un arrondissement…",
                        style="width:100%; box-sizing:border-box; padding:7px 10px; font-size:13px; "
                              "background:#060b14; color:#d9eef8; border:1px solid rgba(0,229,255,0.45); "
                              "border-radius:4px; box-shadow:0 0 10px -3px rgba(0,229,255,0.5);",
                    ),
                    ui.div(id="sidebar-arr-search-results", style="margin-top:2px;"),
                    style="position:relative; margin-bottom:14px;",
                ),
                ui.h6("Calques"),
                ui.input_checkbox("show_bornes", "⚡ Bornes de recharge", True),
                ui.input_checkbox("show_arrond", "🗺️ Arrondissements",    True),
                ui.hr(),
                ui.h6("Coloration des arrondissements"),
                ui.input_selectize(
                    "demo_var", None,
                    choices={
                        "couverture":             "📊 Couverture en bornes",
                        "revenu_median_menage":   "💰 Revenu médian",
                        "tx_voiture_pct":         "🚗 Motorisation (%)",
                        "densite_pop_km2":        "👥 Densité pop. (hab/km²)",
                        "tx_faible_revenu_pct":   "📉 Faible revenu (%)",
                    },
                    selected="couverture",
                    width="100%",
                ),
                ui.hr(),
                ui.p(
                    "Source : Données de Montréal (CC-BY 4.0) · StatCan 2021",
                    style="font-size:0.7rem; color:#999;",
                ),
                width=290,
            ),
            ui.output_ui("map_output"),
        ),
    ),

    # ── TAB 2 : QUESTIONS GESTIONNAIRES ────────────────────────────────────
    ui.nav_panel(
        "🏛 Questions gestionnaires",
        ui.layout_column_wrap(

            # G1 — Parcs
            ui.card(
                ui.card_header("① Parcs — couverture en bornes"),
                ui.p(
                    "Un gestionnaire aimerait savoir si tous les parcs de Montréal "
                    "disposent d'au moins N bornes à 500 m.",
                    style="font-size:0.82rem; color:#555;",
                ),
                ui.layout_columns(
                    ui.input_numeric("g1_n",    "Min. bornes (N)", 20, min=1, max=100),
                    ui.input_action_button("g1_run", "▶ Analyser",
                                           class_="btn btn-primary btn-sm"),
                    col_widths=[7, 5],
                ),
                ui.output_ui("g1_summary"),
                ui.output_table("g1_table"),
            ),

            # G2 — Épiceries
            ui.card(
                ui.card_header("② Épiceries — bornes à proximité"),
                ui.p(
                    "Un gestionnaire souhaite savoir si toutes les épiceries "
                    "de Montréal ont des bornes à 300 m.",
                    style="font-size:0.82rem; color:#555;",
                ),
                ui.input_action_button("g2_run", "▶ Analyser",
                                       class_="btn btn-primary btn-sm"),
                ui.output_ui("g2_summary"),
                ui.output_table("g2_table"),
            ),

            width=1 / 2,
        ),

        ui.br(),

        ui.layout_column_wrap(

            # G3 — Score de priorité
            ui.card(
                ui.card_header("③ Score de priorité composite"),
                ui.p(
                    "Quel arrondissement traiter en premier ? Score 0–100 combinant gap de "
                    "couverture (35%), densité (25%), motorisation (15%), équité de revenu (15%) "
                    "et taux de faible revenu (10%) — 19 arrondissements avec données StatCan.",
                    style="font-size:0.82rem; color:#555;",
                ),
                ui.input_action_button("g3_run", "▶ Calculer les scores",
                                       class_="btn btn-primary btn-sm"),
                ui.output_ui("g3_summary"),
                ui.output_table("g3_table"),
            ),

            # G5 — Intermodalité STM
            ui.card(
                ui.card_header("⑤ Intermodalité STM — park-and-charge"),
                ui.p(
                    "Quelles stations de métro n'ont pas de borne à proximité pour le "
                    "park-and-charge (NOT EXISTS + ST_DWithin, par ligne) ?",
                    style="font-size:0.82rem; color:#555;",
                ),
                ui.layout_columns(
                    ui.input_numeric("g5_rayon", "Rayon (m)", 500, min=100, max=2000, step=100),
                    ui.input_action_button("g5_run", "▶ Analyser",
                                           class_="btn btn-primary btn-sm"),
                    col_widths=[7, 5],
                ),
                ui.output_table("g5_table"),
                ui.output_table("g5_sans_table"),
            ),

            width=1 / 2,
        ),

        ui.br(),

        # G4 — Multi-facteurs
        ui.card(
            ui.card_header("④ Pourquoi cette distribution ? — Analyse multi-facteurs"),
            ui.p(
                "Est-ce lié au profil de la population ? à la proximité de zones de loisirs ? "
                "à la présence de magasins ? — Corrélation Pearson sur 6 facteurs.",
                style="font-size:0.82rem; color:#555;",
            ),
            ui.input_action_button("g6_run", "▶ Analyser les facteurs",
                                   class_="btn btn-primary btn-sm"),
            ui.output_ui("g6_interp"),
            ui.output_table("g6_table"),
        ),
    ),

    # ── TAB 3 : ÉQUITÉ SOCIO-ÉCONOMIQUE ────────────────────────────────────
    ui.nav_panel(
        "📈 Équité socio-éco",
        ui.card(
            ui.card_header(
                "Les bornes dans les quartiers résidentiels sont-elles "
                "conditionnées au profil de la population ?"
            ),
            ui.p(
                "Scatter plot : couverture en bornes (%) vs revenu médian par arrondissement, "
                "avec droite de régression et coefficient de Pearson (r).",
                style="font-size:0.82rem; color:#555;",
            ),
            ui.output_plot("equity_plot", height="480px"),
            ui.output_ui("equity_interp"),
        ),
        ui.br(),
        ui.card(
            ui.card_header("Profil socio-démographique par arrondissement"),
            ui.output_table("equity_table"),
        ),
    ),

    # ── TAB 4 : RAPPORTS ────────────────────────────────────────────────────
    ui.nav_panel(
        "📋 Rapports",
        ui.card(
            ui.card_header("Journal d'activité — recherches et analyses"),
            ui.p(
                "Chaque recherche faite dans la carte et chaque analyse lancée dans "
                "« Questions gestionnaires » est enregistrée ici, sur le disque — le "
                "journal reste disponible même après avoir fermé et relancé l'application.",
                style="font-size:0.82rem; color:#555;",
            ),
            ui.layout_columns(
                ui.input_select(
                    "rapport_type", "Type de rapport",
                    choices={
                        "recherches": "🔎 Historique des recherches (carte)",
                        "analyses":   "📊 Résumé des analyses gestionnaires (G1–G5)",
                        "complet":    "📅 Journal quotidien complet (tout)",
                    },
                    selected="complet",
                ),
                ui.div(
                    ui.input_date("rapport_date", "Date", value=None, language="fr", width="100%"),
                    ui.tags.a(
                        "Toutes les dates", href="#", id="rapport-date-reset-link",
                        style="font-size:0.78rem;",
                    ),
                ),
                col_widths=[7, 5],
            ),
            ui.input_checkbox_group(
                "rapport_categories", "Affiner par outil précis (décoche ce que tu ne veux pas voir)",
                choices={
                    "Recherche carte":      "🔎 Recherche carte",
                    "G1 - Parcs":           "① Parcs",
                    "G2 - Épiceries":       "② Épiceries",
                    "G3 - Score priorité":  "③ Score priorité",
                    "G4 - Corrélation":     "④ Corrélation",
                    "G5 - Intermodalité STM": "⑤ Intermodalité STM",
                },
                selected=["Recherche carte", "G1 - Parcs", "G2 - Épiceries",
                          "G3 - Score priorité", "G4 - Corrélation", "G5 - Intermodalité STM"],
                inline=True,
            ),
            ui.output_ui("rapport_resume"),
            ui.output_table("rapport_table"),
            ui.layout_columns(
                ui.download_button("rapport_download_csv", "⬇ CSV", class_="btn btn-primary btn-sm"),
                ui.download_button("rapport_download_word", "⬇ Word (.docx)", class_="btn btn-primary btn-sm"),
                ui.download_button("rapport_download_pptx", "⬇ PowerPoint (.pptx)", class_="btn btn-primary btn-sm"),
                col_widths=[4, 4, 4],
            ),
        ),
    ),

    # ── TITRE APP ──────────────────────────────────────────────────────────
    title=ui.span(
        "⚡ GeoCharge Montréal",
        style="font-weight:700; letter-spacing:0.5px;",
    ),
    bg="#0d1b2e",
    inverse=True,
    footer=ui.div(
        ui.span(
            ui.img(src=LOGO_DATA_URI, style="height:22px; vertical-align:middle; display:block;"),
            style="display:inline-block; background:#ffffff; border-radius:4px; padding:3px 8px; "
                  "vertical-align:middle; margin-right:10px;",
        ) if LOGO_DATA_URI else None,
        ui.span(
            "GMQ580 — Géomatique Informatique 2 — Été 2026 — Université de Sherbrooke · "
            "Données : Ville de Montréal (CC-BY 4.0) · StatCan Recensement 2021",
            style="vertical-align:middle;",
        ),
        id="app-footer",
        style="text-align:center; font-size:0.7rem; color:#999; padding:8px 4px;",
    ),
)


# ═══════════════════════════════════════════════════════════════════════════
# EXPORT DES RAPPORTS — Word et PowerPoint
# Présentation volontairement sobre (document professionnel classique, pas le
# thème néon de l'app) : un accent de couleur, pas plus.
# ═══════════════════════════════════════════════════════════════════════════
_ACCENT_DOC = DocxRGB(0x00, 0x72, 0x8C)
_ACCENT_PPT = PptxRGB(0x00, 0x72, 0x8C)
_INK_PPT    = PptxRGB(0x16, 0x28, 0x3D)


def _rapport_kpis(df):
    """Indicateurs de synthèse affichés en cases en haut du rapport, dans l'esprit d'un
    tableau de bord (total, catégories distinctes, catégorie dominante, jours actifs)."""
    if df.empty:
        return [
            ("Total événements", "0"),
            ("Catégories", "0"),
            ("Catégorie dominante", "—"),
            ("Jours actifs", "0"),
        ]
    top = df["categorie"].value_counts().idxmax()
    return [
        ("Total événements", f"{len(df):,}".replace(",", " ")),
        ("Catégories", str(df["categorie"].nunique())),
        ("Catégorie dominante", top),
        ("Jours actifs", str(df["date"].nunique())),
    ]


def _rapport_chart_png(df):
    """Répartition des événements par catégorie (barres horizontales), en PNG — sobre,
    même accent de couleur que le reste du rapport."""
    counts = df["categorie"].value_counts().sort_values()
    fig, ax = plt.subplots(figsize=(5.4, 2.8), dpi=160)
    ax.barh(counts.index, counts.values, color="#00728c")
    ax.set_xlabel("Nombre d'événements", fontsize=9, color="#333333")
    ax.tick_params(labelsize=8.5, colors="#333333")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color("#cccccc")
    fig.patch.set_facecolor("white")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _docx_cell_bg(cell, hex_color):
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    cell._tc.get_or_add_tcPr().append(shd)


def _build_rapport_docx(df, description_filtre):
    doc = docx.Document()

    if os.path.exists(LOGO_PATH):
        p_logo = doc.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p_logo.add_run().add_picture(LOGO_PATH, width=Cm(3.2))

    titre = doc.add_paragraph()
    r = titre.add_run("GeoCharge Montréal — Rapport d'activité")
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = _ACCENT_DOC

    sous_titre = doc.add_paragraph()
    r = sous_titre.add_run(
        f"Généré le {datetime.now().strftime('%Y-%m-%d à %H:%M')} — {description_filtre}"
    )
    r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = DocxRGB(0x60, 0x60, 0x60)

    doc.add_paragraph()

    # ── Cases KPI (synthèse en un coup d'œil, avant le détail du journal) ────
    kpis = _rapport_kpis(df)
    kpi_table = doc.add_table(rows=2, cols=len(kpis))
    kpi_table.autofit = True
    for i, (label, valeur) in enumerate(kpis):
        cell_valeur = kpi_table.cell(0, i)
        _docx_cell_bg(cell_valeur, "E6F2F5")
        p_val = cell_valeur.paragraphs[0]
        p_val.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_val = p_val.add_run(valeur)
        r_val.font.size = Pt(18); r_val.font.bold = True; r_val.font.color.rgb = _ACCENT_DOC

        cell_label = kpi_table.cell(1, i)
        _docx_cell_bg(cell_label, "E6F2F5")
        p_lab = cell_label.paragraphs[0]
        p_lab.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_lab = p_lab.add_run(label)
        r_lab.font.size = Pt(8.5); r_lab.font.color.rgb = DocxRGB(0x50, 0x50, 0x50)

    doc.add_paragraph()

    # ── Répartition par catégorie (graphique) ────────────────────────────────
    if not df.empty:
        p_chart = doc.add_paragraph()
        p_chart.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_chart.add_run().add_picture(_rapport_chart_png(df), width=Cm(14))
        doc.add_paragraph()

    if df.empty:
        p = doc.add_paragraph("Aucun événement pour ce filtre.")
        p.runs[0].font.italic = True
    else:
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        for i, col in enumerate(["Date", "Heure", "Catégorie", "Détail"]):
            hdr[i].text = col
            hdr[i].paragraphs[0].runs[0].font.bold = True
        for _, row in df.iterrows():
            cells = table.add_row().cells
            cells[0].text = str(row["date"])
            cells[1].text = str(row["heure"])
            cells[2].text = str(row["categorie"])
            cells[3].text = str(row["detail"])
        table.columns[0].width = Cm(2.4)
        table.columns[1].width = Cm(2.0)
        table.columns[2].width = Cm(3.6)
        table.columns[3].width = Cm(8.0)

    doc.add_paragraph()
    pied = doc.add_paragraph()
    r = pied.add_run("GMQ580 — Géomatique Informatique 2 — Été 2026 — Université de Sherbrooke")
    r.font.size = Pt(8); r.font.color.rgb = DocxRGB(0x90, 0x90, 0x90)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_rapport_pptx(df, description_filtre):
    prs = Presentation()
    prs.slide_width = PptxIn(13.333)
    prs.slide_height = PptxIn(7.5)
    blank = prs.slide_layouts[6]

    # ── Diapositive de titre (sobre) ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    if os.path.exists(LOGO_PATH):
        s.shapes.add_picture(LOGO_PATH, PptxIn(11.3), PptxIn(0.4), height=PptxIn(0.7))
    tb = s.shapes.add_textbox(PptxIn(0.7), PptxIn(2.6), PptxIn(12), PptxIn(1.2))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run(); run.text = "GeoCharge Montréal — Rapport d'activité"
    run.font.size = PptxPt(32); run.font.bold = True; run.font.color.rgb = _INK_PPT

    tb2 = s.shapes.add_textbox(PptxIn(0.7), PptxIn(3.5), PptxIn(12), PptxIn(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"Généré le {datetime.now().strftime('%Y-%m-%d à %H:%M')} — {description_filtre}"
    run2.font.size = PptxPt(14); run2.font.italic = True; run2.font.color.rgb = PptxRGB(0x70, 0x70, 0x70)

    line = s.shapes.add_shape(1, PptxIn(0.7), PptxIn(2.45), PptxIn(3.0), PptxPt(3))
    line.fill.solid(); line.fill.fore_color.rgb = _ACCENT_PPT; line.line.fill.background()

    # ── Diapositive tableau de bord : cases KPI + répartition par catégorie ──
    s_dash = prs.slides.add_slide(blank)
    tb_dash = s_dash.shapes.add_textbox(PptxIn(0.6), PptxIn(0.35), PptxIn(12), PptxIn(0.5))
    run_dash = tb_dash.text_frame.paragraphs[0].add_run()
    run_dash.text = "Synthèse"
    run_dash.font.size = PptxPt(20); run_dash.font.bold = True; run_dash.font.color.rgb = _ACCENT_PPT

    kpis = _rapport_kpis(df)
    card_w = 2.9
    gap = 0.25
    start_x = (13.333 - (card_w * len(kpis) + gap * (len(kpis) - 1))) / 2
    for i, (label, valeur) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        card = s_dash.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, PptxIn(x), PptxIn(1.1), PptxIn(card_w), PptxIn(1.3)
        )
        card.fill.solid(); card.fill.fore_color.rgb = PptxRGB(0xE6, 0xF2, 0xF5)
        card.line.color.rgb = _ACCENT_PPT; card.line.width = PptxPt(0.75)
        tf = card.text_frame
        tf.word_wrap = True
        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        r_val = p_val.add_run(); r_val.text = valeur
        r_val.font.size = PptxPt(22); r_val.font.bold = True; r_val.font.color.rgb = _ACCENT_PPT
        p_lab = tf.add_paragraph()
        p_lab.alignment = PP_ALIGN.CENTER
        r_lab = p_lab.add_run(); r_lab.text = label
        r_lab.font.size = PptxPt(10); r_lab.font.color.rgb = PptxRGB(0x50, 0x50, 0x50)

    if not df.empty:
        s_dash.shapes.add_picture(
            _rapport_chart_png(df), PptxIn(2.1), PptxIn(2.8), width=PptxIn(9.1)
        )
    else:
        tb_empty = s_dash.shapes.add_textbox(PptxIn(0.7), PptxIn(3.3), PptxIn(12), PptxIn(1))
        run_empty = tb_empty.text_frame.paragraphs[0].add_run()
        run_empty.text = "Aucun événement pour ce filtre."
        run_empty.font.size = PptxPt(16); run_empty.font.italic = True
        run_empty.font.color.rgb = PptxRGB(0x70, 0x70, 0x70)

    # ── Diapositive(s) de tableau — 16 lignes par diapo pour rester lisible ──
    PAR_PAGE = 16
    lignes = df.to_dict("records") if not df.empty else []
    if not lignes:
        s2 = prs.slides.add_slide(blank)
        tb = s2.shapes.add_textbox(PptxIn(0.7), PptxIn(3.3), PptxIn(12), PptxIn(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Aucun événement pour ce filtre."
        run.font.size = PptxPt(18); run.font.italic = True; run.font.color.rgb = PptxRGB(0x70, 0x70, 0x70)

    for debut in range(0, len(lignes), PAR_PAGE):
        page = lignes[debut:debut + PAR_PAGE]
        s2 = prs.slides.add_slide(blank)
        tb = s2.shapes.add_textbox(PptxIn(0.6), PptxIn(0.35), PptxIn(12), PptxIn(0.5))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = f"Journal d'activité ({debut + 1}–{debut + len(page)} sur {len(lignes)})"
        run.font.size = PptxPt(16); run.font.bold = True; run.font.color.rgb = _ACCENT_PPT

        rows, cols = len(page) + 1, 4
        tbl_shape = s2.shapes.add_table(rows, cols, PptxIn(0.6), PptxIn(1.0), PptxIn(12.1), PptxIn(6.1))
        tbl = tbl_shape.table
        tbl.columns[0].width = PptxIn(1.6); tbl.columns[1].width = PptxIn(1.3)
        tbl.columns[2].width = PptxIn(2.6); tbl.columns[3].width = PptxIn(6.6)
        for i, col in enumerate(["Date", "Heure", "Catégorie", "Détail"]):
            c = tbl.cell(0, i); c.text = col
            c.text_frame.paragraphs[0].font.size = PptxPt(11)
            c.text_frame.paragraphs[0].font.bold = True
            c.fill.solid(); c.fill.fore_color.rgb = _ACCENT_PPT
            c.text_frame.paragraphs[0].font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)
        for r_i, row in enumerate(page, start=1):
            vals = [str(row["date"]), str(row["heure"]), str(row["categorie"]), str(row["detail"])]
            for c_i, v in enumerate(vals):
                c = tbl.cell(r_i, c_i); c.text = v
                c.text_frame.paragraphs[0].font.size = PptxPt(10)
                c.fill.solid(); c.fill.fore_color.rgb = PptxRGB(0xFF, 0xFF, 0xFF) if r_i % 2 else PptxRGB(0xF2, 0xF6, 0xF8)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# SERVEUR SHINY
# ═══════════════════════════════════════════════════════════════════════════

def server(input, output, session):

    # ── JOURNAL D'ACTIVITÉ (persistant sur disque) ──────────────────────────
    journal_rv = reactive.value(_lire_journal())

    def _log_and_refresh(categorie, detail):
        _log_event(categorie, detail)
        journal_rv.set(_lire_journal())

    @reactive.effect
    def _update_rapport_dates():
        df = journal_rv()
        # Restreint le calendrier à la période couverte par le journal, plutôt qu'une liste
        # de dates à choisir dans un menu déroulant — le calendrier grise le reste tout seul.
        if df.empty:
            return
        ui.update_date("rapport_date", min=df["date"].min(), max=df["date"].max())

    # ── CARTE ──────────────────────────────────────────────────────────────
    @render.ui
    def map_output():
        m = build_map(
            show_bornes=input.show_bornes(),
            show_arrond=input.show_arrond(),
            demo_var=input.demo_var(),
        )
        html = m._repr_html_()
        # folium enveloppe l'iframe dans un div "aspect-ratio" (height:0 + padding-bottom:%),
        # pas dans un iframe à hauteur fixe. C'est ce div qu'il faut cibler pour que l'iframe
        # (position:absolute; height:100%) remplisse l'espace réellement disponible, ajusté
        # dynamiquement en JS (voir ADJUST_MAP_HEIGHT_JS) plutôt qu'un calc(100vh - Xpx) figé.
        html = html.replace(
            'style="position:relative;width:100%;height:0;padding-bottom:60%;"',
            'id="geocharge-map-frame" style="position:relative;width:100%;height:70vh;"',
        )
        return ui.HTML(html)

    @reactive.effect
    @reactive.event(input.arr_search_click)
    def _log_search():
        nom = input.arr_search_click().get("nom", "?")
        _log_and_refresh("Recherche carte", f"Arrondissement recherché : {nom}")

    # ── G1 : PARCS ─────────────────────────────────────────────────────────
    @render.ui
    @reactive.event(input.g1_run)
    def g1_summary():
        n     = input.g1_n()
        insuf = parcs_gdf[parcs_gdf["nb_bornes_500m"] < n]
        total = len(parcs_gdf)
        count = len(insuf)
        pct_ok = round(100 * (total - count) / max(total, 1), 1)
        couleur = "success" if pct_ok >= 80 else "warning" if pct_ok >= 50 else "danger"
        return ui.div(
            ui.div(
                ui.strong(f"{count}"),
                f" parcs sur {total} ont moins de {n} bornes dans un rayon de 500 m.",
                class_=f"alert alert-{couleur}",
                style="font-size:0.88rem; padding:8px 12px; margin:8px 0;",
            ),
            ui.p(
                f"✅ {pct_ok}% des parcs atteignent le seuil cible de {n} bornes à 500 m.",
                style="font-size:0.82rem; color:#555;",
            ),
        )

    @render.table
    @reactive.event(input.g1_run)
    def g1_table():
        n     = input.g1_n()
        insuf = parcs_gdf[parcs_gdf["nb_bornes_500m"] < n].copy()
        insuf = insuf.sort_values("nb_bornes_500m")[["nom", "superficie_ha", "nb_bornes_500m"]].head(25)
        insuf.columns = ["Parc", "Superficie (ha)", f"Bornes à 500 m (< {n})"]
        return insuf.reset_index(drop=True)

    @reactive.effect
    @reactive.event(input.g1_run)
    def _log_g1():
        n = input.g1_n()
        count = int((parcs_gdf["nb_bornes_500m"] < n).sum())
        _log_and_refresh("G1 - Parcs", f"Seuil N={n} bornes à 500 m → {count}/{len(parcs_gdf)} parcs sous le seuil")

    # ── G2 : ÉPICERIES ─────────────────────────────────────────────────────
    @render.ui
    @reactive.event(input.g2_run)
    def g2_summary():
        sans   = epiceries_gdf[~epiceries_gdf["has_borne_300m"]]
        total  = len(epiceries_gdf)
        count  = len(sans)
        pct_ok = round(100 * (total - count) / max(total, 1), 1)
        couleur = "success" if pct_ok >= 80 else "warning" if pct_ok >= 50 else "danger"
        return ui.div(
            ui.div(
                ui.strong(f"{count}"),
                f" épiceries sur {total} n'ont aucune borne dans un rayon de 300 m.",
                class_=f"alert alert-{couleur}",
                style="font-size:0.88rem; padding:8px 12px; margin:8px 0;",
            ),
            ui.p(
                f"✅ {pct_ok}% des épiceries ont au moins une borne à 300 m.",
                style="font-size:0.82rem; color:#555;",
            ),
        )

    @render.table
    @reactive.event(input.g2_run)
    def g2_table():
        sans = epiceries_gdf[~epiceries_gdf["has_borne_300m"]][["nom", "type", "adresse"]].head(25).copy()
        sans.columns = ["Épicerie", "Type", "Adresse"]
        return sans.reset_index(drop=True)

    @reactive.effect
    @reactive.event(input.g2_run)
    def _log_g2():
        count = int((~epiceries_gdf["has_borne_300m"]).sum())
        _log_and_refresh("G2 - Épiceries", f"{count}/{len(epiceries_gdf)} épiceries sans borne à 300 m")

    # ── G3 : SCORE DE PRIORITÉ ─────────────────────────────────────────────
    @render.ui
    @reactive.event(input.g3_run)
    def g3_summary():
        top = priorite_df.iloc[0]
        return ui.div(
            ui.strong(f"{top['nom']}"),
            f" arrive en tête avec un score de {top['score_priorite']}/100 "
            f"(couverture actuelle : {top['pct_couverture']}%).",
            class_="alert alert-info",
            style="font-size:0.88rem; padding:8px 12px; margin:8px 0;",
        )

    @render.table
    @reactive.event(input.g3_run)
    def g3_table():
        df = priorite_df.head(10).copy()
        df.columns = ["Arrondissement", "Couverture (%)", "Nb bornes", "Score priorité /100"]
        return df

    @reactive.effect
    @reactive.event(input.g3_run)
    def _log_g3():
        top = priorite_df.iloc[0]
        _log_and_refresh("G3 - Score priorité",
                          f"1er : {top['nom']} ({top['score_priorite']}/100)")

    # ── G5 : INTERMODALITÉ STM ─────────────────────────────────────────────
    @reactive.calc
    @reactive.event(input.g5_run)
    def _g5_result():
        rayon = input.g5_rayon()
        buf_union_r = bornes_p.geometry.buffer(rayon).unary_union
        stations = stations_gdf.copy()
        stations["a_borne"] = stations_p.geometry.within(buf_union_r)

        rows = []
        for ligne in ["Verte", "Orange", "Bleue", "Jaune"]:
            sub = stations[stations["ligne"] == ligne]
            total = len(sub)
            avec = int(sub["a_borne"].sum())
            rows.append({
                "Ligne": ligne, "Stations": total, "Avec borne": avec,
                "Sans borne": total - avec,
                "Couverture (%)": round(100 * avec / total, 1) if total else 0.0,
            })
        total_all = len(stations)
        avec_all  = int(stations["a_borne"].sum())
        rows.append({
            "Ligne": "Total", "Stations": total_all, "Avec borne": avec_all,
            "Sans borne": total_all - avec_all,
            "Couverture (%)": round(100 * avec_all / total_all, 1) if total_all else 0.0,
        })
        return pd.DataFrame(rows), stations[~stations["a_borne"]][["nom", "ligne"]]

    @render.table
    @reactive.event(input.g5_run)
    def g5_table():
        return _g5_result()[0]

    @render.table
    @reactive.event(input.g5_run)
    def g5_sans_table():
        df = _g5_result()[1].copy().reset_index(drop=True)
        df.columns = ["Station sans borne à proximité", "Ligne"]
        return df

    @reactive.effect
    @reactive.event(input.g5_run)
    def _log_g5():
        rayon = input.g5_rayon()
        total_row, _ = _g5_result()
        row = total_row[total_row["Ligne"] == "Total"].iloc[0]
        _log_and_refresh("G5 - Intermodalité STM",
                          f"Rayon {rayon} m → {row['Sans borne']}/{row['Stations']} stations sans borne")

    # ── G6 : MULTI-FACTEURS ────────────────────────────────────────────────
    @reactive.calc
    @reactive.event(input.g6_run)
    def _g6_cors():
        cors = {}
        for var in VAR_LABELS_G6:
            if var not in g6_df.columns:
                continue
            sub = g6_df.dropna(subset=[var, "pct_couverture"])
            cors[var] = pearson(sub[var].tolist(), sub["pct_couverture"].tolist())
        return sorted(
            [{"variable": k, "r": v,
              "label": VAR_LABELS_G6[k],
              "question": QUESTIONS_G6.get(k, "")} for k, v in cors.items()],
            key=lambda x: -abs(x["r"])
        )

    @render.ui
    @reactive.event(input.g6_run)
    def g6_interp():
        cors = _g6_cors()
        if not cors:
            return ui.p("Données insuffisantes.")
        top = cors[0]
        absR = abs(top["r"])
        if absR >= 0.4:
            direction = "positive" if top["r"] > 0 else "négative"
            txt = (f"Facteur dominant : « {top['label']} » (r = {top['r']}) — "
                   f"corrélation {direction} forte. Les zones avec "
                   f"{'plus de' if top['r']>0 else 'moins de'} {top['label'].lower()} "
                   f"tendent à avoir une meilleure couverture en bornes.")
            cls = "alert alert-info"
        elif absR >= 0.2:
            txt = (f"Facteur le plus lié : « {top['label']} » (r = {top['r']}) — "
                   f"corrélation modérée. La distribution est probablement multifactorielle.")
            cls = "alert alert-warning"
        else:
            txt = ("Aucun facteur isolé n'explique clairement la distribution. "
                   "La localisation des bornes est davantage liée à des décisions historiques "
                   "ou à la disponibilité foncière qu'à un critère socio-spatial unique.")
            cls = "alert alert-secondary"
        return ui.div(txt, class_=cls, style="font-size:0.88rem; margin:8px 0;")

    @render.table
    @reactive.event(input.g6_run)
    def g6_table():
        cors = _g6_cors()
        rows = []
        for f in cors:
            absR = abs(f["r"])
            force = "Forte" if absR >= 0.4 else "Modérée" if absR >= 0.2 else "Faible"
            dir_  = "↑ positive" if f["r"] > 0.05 else ("↓ négative" if f["r"] < -0.05 else "≈ nulle")
            rows.append({
                "Facteur":       f["label"],
                "r (Pearson)":   f["r"],
                "Direction":     dir_,
                "Force":         force,
                "Question prof": f["question"],
            })
        return pd.DataFrame(rows).reset_index(drop=True)

    @reactive.effect
    @reactive.event(input.g6_run)
    def _log_g6():
        cors = _g6_cors()
        if cors:
            top = cors[0]
            _log_and_refresh("G4 - Corrélation", f"Facteur dominant : {top['label']} (r={top['r']})")

    # ── ÉQUITÉ ─────────────────────────────────────────────────────────────
    @render.plot
    def equity_plot():
        df = arrond_gdf[["nom", "pct_couverture",
                          "revenu_median_menage", "tx_faible_revenu_pct"]].dropna()

        fig, ax = plt.subplots(figsize=(11, 6))
        fig.patch.set_facecolor("#0d1b2e")
        ax.set_facecolor("#0a0e1a")

        sc = ax.scatter(
            df["revenu_median_menage"] / 1000,
            df["pct_couverture"],
            c=df["tx_faible_revenu_pct"],
            cmap="YlOrRd_r",
            s=90, alpha=0.9,
            edgecolors="#00d4ff", linewidths=0.6,
        )

        for _, row in df.iterrows():
            ax.annotate(
                row["nom"].split("–")[0].strip(),
                (row["revenu_median_menage"] / 1000, row["pct_couverture"]),
                fontsize=5.5, color="#9ab4cc", ha="center", va="bottom",
                xytext=(0, 4), textcoords="offset points",
            )

        # Droite de régression
        xs = df["revenu_median_menage"].values
        ys = df["pct_couverture"].values
        coeffs = np.polyfit(xs, ys, 1)
        x_line = np.linspace(xs.min(), xs.max(), 100)
        ax.plot(x_line / 1000, np.polyval(coeffs, x_line),
                color="#00d4ff", linewidth=1.5, linestyle="--", alpha=0.7,
                label="Régression linéaire")

        r = pearson(xs.tolist(), ys.tolist())
        ax.set_xlabel("Revenu médian des ménages (k$)",  color="#dce8f4", fontsize=10)
        ax.set_ylabel("Couverture en bornes (%)",         color="#dce8f4", fontsize=10)
        ax.set_title(
            f"Équité socio-économique — couverture vs revenu  (r = {r})\n"
            "Couleur = taux de faible revenu (rouge = vulnérable)",
            color="#dce8f4", fontsize=11,
        )
        ax.tick_params(colors="#6a8aaa", labelsize=8)
        for sp in ax.spines.values():
            sp.set_color("#1e3a5f")
        ax.legend(fontsize=8, facecolor="#0d1b2e", labelcolor="#dce8f4",
                  framealpha=0.7)

        cbar = plt.colorbar(sc, ax=ax)
        cbar.set_label("Faible revenu (%)", color="#dce8f4", fontsize=8)
        cbar.ax.yaxis.set_tick_params(color="#6a8aaa", labelcolor="#6a8aaa")

        plt.tight_layout(pad=1.5)
        return fig

    @render.ui
    def equity_interp():
        df = arrond_gdf[["pct_couverture", "revenu_median_menage"]].dropna()
        r = pearson(df["revenu_median_menage"].tolist(), df["pct_couverture"].tolist())
        if r > 0.25:
            txt = (f"r = {r} — inéquité détectée : les arrondissements à revenu élevé "
                   "tendent à être mieux desservis en bornes.")
            cls = "alert alert-warning"
        elif r < -0.25:
            txt = (f"r = {r} — pas d'inéquité en faveur des riches : à Montréal, la couverture "
                   "est en fait légèrement plus faible dans les arrondissements à revenu élevé "
                   "(souvent plus périphériques et moins denses). Le revenu n'est donc pas le "
                   "facteur qui explique le mieux la distribution des bornes.")
            cls = "alert alert-info"
        else:
            txt = (f"r = {r} — équité relative : pas de corrélation significative entre "
                   "revenu médian et couverture en bornes.")
            cls = "alert alert-secondary"
        return ui.div(txt, class_=cls, style="font-size:0.88rem; margin:8px 0;")

    @render.table
    def equity_table():
        df = arrond_gdf[["nom", "pct_couverture", "nb_bornes",
                          "revenu_median_menage", "tx_voiture_pct",
                          "tx_faible_revenu_pct", "densite_pop_km2"]].copy()
        df = df.sort_values("pct_couverture").rename(columns={
            "nom":                  "Arrondissement",
            "pct_couverture":       "Couverture (%)",
            "nb_bornes":            "Nb bornes",
            "revenu_median_menage": "Revenu médian ($)",
            "tx_voiture_pct":       "Motorisation (%)",
            "tx_faible_revenu_pct": "Faible revenu (%)",
            "densite_pop_km2":      "Densité (hab/km²)",
        })
        return df.reset_index(drop=True)

    # ── RAPPORTS ───────────────────────────────────────────────────────────
    CATEGORIES_RECHERCHE = ["Recherche carte"]
    CATEGORIES_ANALYSES  = ["G1 - Parcs", "G2 - Épiceries", "G3 - Score priorité",
                             "G4 - Corrélation", "G5 - Intermodalité STM"]

    @reactive.calc
    def _rapport_filtre():
        df = journal_rv()
        type_ = input.rapport_type()
        if type_ == "recherches":
            df = df[df["categorie"].isin(CATEGORIES_RECHERCHE)]
        elif type_ == "analyses":
            df = df[df["categorie"].isin(CATEGORIES_ANALYSES)]
        date_ = input.rapport_date()
        if date_ is not None:
            df = df[df["date"] == str(date_)]
        cats_precises = input.rapport_categories()
        if cats_precises:
            df = df[df["categorie"].isin(cats_precises)]
        else:
            df = df.iloc[0:0]
        return df.sort_values(["date", "heure"], ascending=False).reset_index(drop=True)

    @render.ui
    def rapport_resume():
        df = _rapport_filtre()
        if df.empty:
            return ui.div(
                "Aucun événement pour ce filtre — vérifie que tu n'as pas décoché toutes "
                "les catégories, ou fais une recherche sur la carte / lance une analyse "
                "dans « Questions gestionnaires » pour en voir apparaître ici.",
                class_="alert alert-secondary",
                style="font-size:0.85rem; padding:8px 12px; margin:8px 0;",
            )
        troncature = " (affichage limité aux 200 plus récents)" if len(df) > 200 else ""
        return ui.div(
            ui.strong(f"{len(df)}"), f" événement(s){troncature} — journal complet : ",
            ui.code(JOURNAL_PATH),
            class_="alert alert-info",
            style="font-size:0.82rem; padding:8px 12px; margin:8px 0;",
        )

    @render.table
    def rapport_table():
        df = _rapport_filtre().copy()
        df.columns = ["Date", "Heure", "Catégorie", "Détail"]
        return df.head(200)

    def _description_filtre():
        type_labels = {"recherches": "Historique des recherches", "analyses": "Résumé des analyses",
                       "complet": "Journal quotidien complet"}
        n_cat = len(input.rapport_categories())
        date_ = input.rapport_date()
        date_txt = "toutes les dates" if date_ is None else str(date_)
        return (f"{type_labels.get(input.rapport_type(), 'Journal')} · {date_txt} · "
                f"{n_cat}/6 catégories affichées")

    _NOM_FICHIER = lambda ext: f"rapport_geocharge_{datetime.now().strftime('%Y%m%d_%H%M')}.{ext}"

    @render.download(filename=lambda: _NOM_FICHIER("csv"))
    def rapport_download_csv():
        buf = io.StringIO()
        _rapport_filtre().to_csv(buf, index=False)
        yield buf.getvalue()

    @render.download(filename=lambda: _NOM_FICHIER("docx"))
    def rapport_download_word():
        yield _build_rapport_docx(_rapport_filtre(), _description_filtre())

    @render.download(filename=lambda: _NOM_FICHIER("pptx"))
    def rapport_download_pptx():
        yield _build_rapport_pptx(_rapport_filtre(), _description_filtre())


# ═══════════════════════════════════════════════════════════════════════════
app = App(app_ui, server)
